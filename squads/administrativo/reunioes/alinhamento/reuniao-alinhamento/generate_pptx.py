from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
AZUL_ESCURO = RGBColor(0x00, 0x33, 0x66)
AZUL_MEDIO = RGBColor(0x00, 0x66, 0xCC)
AZUL_CLARO = RGBColor(0x00, 0x99, 0xFF)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_FUNDO = RGBColor(0xF0, 0xF4, 0xF8)
CINZA_TEXTO = RGBColor(0x4A, 0x55, 0x68)
VERDE = RGBColor(0x28, 0xA7, 0x45)
VERMELHO = RGBColor(0xDC, 0x35, 0x45)
AMARELO = RGBColor(0xFF, 0xC1, 0x07)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BRANCO):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=AZUL_ESCURO, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=16, bold=False, color=AZUL_ESCURO, alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name='Calibri'):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p

def add_bullet(text_frame, text, font_size=16, color=AZUL_ESCURO, bold=False, level=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.font.bold = bold
    p.level = level
    p.space_before = Pt(4)
    return p

def make_king_number(slide, number, label, left, top, width, height):
    add_textbox(slide, left, top, width, Inches(1), number, font_size=48, bold=True, color=AZUL_MEDIO, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(1), width, Inches(0.5), label, font_size=14, color=CINZA_TEXTO, alignment=PP_ALIGN.CENTER)

def make_slide_number(slide, num):
    add_textbox(slide, Inches(12.3), Inches(7.0), Inches(1), Inches(0.4), str(num), font_size=11, color=CINZA_TEXTO, alignment=PP_ALIGN.RIGHT)

def make_header_bar(slide, text):
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), AZUL_ESCURO)
    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8), text, font_size=28, bold=True, color=BRANCO)

# ===================== SLIDE 1: CAPA =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRANCO)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(3.5), AZUL_ESCURO)
add_textbox(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.2), 'REUNIÃO DE ALINHAMENTO', font_size=40, bold=True, color=BRANCO, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1), 'Maio / 2026', font_size=32, color=AZUL_CLARO, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.0), Inches(11), Inches(1), 'ProSystem Sistemas', font_size=22, color=AZUL_ESCURO, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6), 'Equipe ProSystem', font_size=16, color=CINZA_TEXTO, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.5), Inches(3.5), Inches(2.333), Inches(0.06), AZUL_MEDIO)

# ===================== SLIDE 2: AGENDA =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRANCO)
make_header_bar(slide, 'PAUTA DO DIA')

items = [
    ('🛠️', 'Padrão de Atendimento Técnico'),
    ('📢', 'Marketing & Presença Digital'),
    ('💼', 'Negociações em Andamento'),
    ('📊', 'Clientes Perdidos — Análise'),
    ('⭐', 'Avaliações Google & Meta 100'),
    ('📌', 'Encaminhamentos'),
]

y_start = 1.8
for i, (icon, text) in enumerate(items):
    y = y_start + i * 0.85
    add_shape(slide, Inches(1.5), Inches(y), Inches(0.6), Inches(0.6), AZUL_MEDIO, alpha=1000)
    add_textbox(slide, Inches(1.55), Inches(y + 0.05), Inches(0.5), Inches(0.5), icon, font_size=20, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.4), Inches(y + 0.1), Inches(9), Inches(0.5), text, font_size=22, color=AZUL_ESCURO)

make_slide_number(slide, 2)

# ===================== SLIDE 3: ATENDIMENTO TÉCNICO =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CINZA_FUNDO)
make_header_bar(slide, 'PADRÃO DE ATENDIMENTO TÉCNICO')

cards = [
    ('INÍCIO → MEIO → FIM', 'Ciclo completo até a resolução final do cliente'),
    ('USO DA AGENDA', 'Organizar, acompanhar e finalizar cada demanda'),
    ('PROATIVIDADE', 'Não esperar o cliente chamar — tomar a iniciativa'),
    ('ATIVO LIBERADO', 'Agenda livre para foco total nas demandas ativas'),
]

for i, (title, desc) in enumerate(cards):
    x = 1.0 + i * 3.0
    card = add_shape(slide, Inches(x), Inches(1.8), Inches(2.7), Inches(3.5), BRANCO)
    card.shadow.inherit = False
    add_shape(slide, Inches(x), Inches(1.8), Inches(2.7), Inches(0.15), AZUL_MEDIO)
    add_textbox(slide, Inches(x + 0.2), Inches(2.2), Inches(2.3), Inches(1), title, font_size=16, bold=True, color=AZUL_ESCURO, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x + 0.2), Inches(3.2), Inches(2.3), Inches(1.5), desc, font_size=14, color=CINZA_TEXTO, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.6), '🎯  "Agenda livre = atenção total ao cliente"', font_size=18, bold=True, color=AZUL_MEDIO, alignment=PP_ALIGN.CENTER)
make_slide_number(slide, 3)

# ===================== SLIDE 4: MARKETING =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRANCO)
make_header_bar(slide, 'MARKETING — MAIO/2026')

marketing_items = [
    ('🌐', 'Site ProSystem', 'prosystemnet.com — renovado com soluções completas'),
    ('📝', 'Blog Semanal', 'Postagens automáticas para engajamento e visibilidade'),
    ('📢', 'Campanhas Ativas', 'Farmácias + Padarias — leads chegando diariamente'),
]

for i, (icon, title, desc) in enumerate(marketing_items):
    y = 1.8 + i * 1.6
    add_shape(slide, Inches(1.2), Inches(y), Inches(0.8), Inches(0.8), AZUL_MEDIO)
    add_textbox(slide, Inches(1.25), Inches(y + 0.1), Inches(0.7), Inches(0.7), icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.3), Inches(y), Inches(4), Inches(0.5), title, font_size=22, bold=True, color=AZUL_ESCURO)
    add_textbox(slide, Inches(2.3), Inches(y + 0.5), Inches(9), Inches(0.4), desc, font_size=16, color=CINZA_TEXTO)

add_shape(slide, Inches(1.2), Inches(6.2), Inches(10.5), Inches(0.06), AZUL_CLARO)
add_textbox(slide, Inches(1), Inches(6.4), Inches(11), Inches(0.6), '🎯  Meta: Mais conteúdo → Mais visibilidade → Mais clientes', font_size=18, bold=True, color=AZUL_MEDIO, alignment=PP_ALIGN.CENTER)
make_slide_number(slide, 4)

# ===================== SLIDE 5: GOOGLE REVIEWS =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CINZA_FUNDO)
make_header_bar(slide, 'GOOGLE REVIEWS — RUMO ÀS 100!')

make_king_number(slide, '4,9 ⭐', 'Nota atual', Inches(1.5), Inches(1.8), Inches(3), Inches(1.5))
make_king_number(slide, '87 / 100', 'avaliações', Inches(5.2), Inches(1.8), Inches(3), Inches(1.5))
make_king_number(slide, '+36', 'novas em Maio', Inches(8.9), Inches(1.8), Inches(3), Inches(1.5))

add_shape(slide, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.04), AZUL_CLARO)

add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.4), 'AVALIAÇÕES EM DESTAQUE', font_size=18, bold=True, color=AZUL_ESCURO)
add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.8), '⭐ "Muito bem assessorado... ótimo atendimento e serviços relevantes de suporte"', font_size=14, color=CINZA_TEXTO)
add_textbox(slide, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.4), '   — Arthur, Farmácia Menor Preço (Cujubim-RO)', font_size=13, bold=True, color=AZUL_MEDIO)

add_textbox(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.8), '⭐ "Bom empenho na prestação de serviços de instalação e suporte técnico"', font_size=14, color=CINZA_TEXTO)
add_textbox(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.4), '   — Welder, LOOSE FARMA (Serra-ES)', font_size=13, bold=True, color=AZUL_MEDIO)

add_textbox(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.8), '⭐ "Fácil de lidar, bem objetivo e com atenciosos profissionais sempre a disposição"', font_size=14, color=CINZA_TEXTO)
add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4), '   — Magno Campos, Farmácia da Vila (Conceição da Barra-ES)', font_size=13, bold=True, color=AZUL_MEDIO)

make_slide_number(slide, 5)

# ===================== SLIDE 6: NEGOCIAÇÕES FECHADAS =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRANCO)
make_header_bar(slide, 'NEGOCIACOES FECHADAS')

# King numbers row
make_king_number(slide, 'R$ 24.180', 'caixa imediato (instalacoes)', Inches(0.8), Inches(1.6), Inches(3.8), Inches(1.5))
make_king_number(slide, 'R$ 5.450/mes', 'MMR mensal', Inches(4.8), Inches(1.6), Inches(3.8), Inches(1.5))
make_king_number(slide, 'R$ 1.612', 'media de implantacao', Inches(8.8), Inches(1.6), Inches(3.8), Inches(1.5))

# Cards - Farma Plus e Padaria
add_shape(slide, Inches(0.8), Inches(3.5), Inches(5.8), Inches(2.8), BRANCO)
add_shape(slide, Inches(0.8), Inches(3.5), Inches(5.8), Inches(0.15), AZUL_MEDIO)
add_textbox(slide, Inches(1.2), Inches(3.8), Inches(5), Inches(0.5), 'FARMA PLUS  (13 propostas)', font_size=20, bold=True, color=AZUL_ESCURO)
add_textbox(slide, Inches(1.2), Inches(4.4), Inches(5), Inches(0.3), 'Media mensalidade: R$ 376 / mes', font_size=16, color=CINZA_TEXTO)
add_textbox(slide, Inches(1.2), Inches(4.8), Inches(5), Inches(0.3), 'Media implantacao: R$ 1.752', font_size=16, color=CINZA_TEXTO)
add_textbox(slide, Inches(1.2), Inches(5.3), Inches(5), Inches(0.3), 'Expansao em Goias e Rondonia', font_size=16, bold=True, color=AZUL_MEDIO)

add_shape(slide, Inches(7), Inches(3.5), Inches(5.5), Inches(2.8), BRANCO)
add_shape(slide, Inches(7), Inches(3.5), Inches(5.5), Inches(0.15), VERDE)
add_textbox(slide, Inches(7.4), Inches(3.8), Inches(4.8), Inches(0.5), 'PADARIA  (2 propostas)', font_size=20, bold=True, color=AZUL_ESCURO)
add_textbox(slide, Inches(7.4), Inches(4.4), Inches(4.8), Inches(0.3), 'Media mensalidade: R$ 280 / mes', font_size=16, color=CINZA_TEXTO)
add_textbox(slide, Inches(7.4), Inches(4.8), Inches(4.8), Inches(0.3), 'Media implantacao: R$ 700', font_size=16, color=CINZA_TEXTO)

# Status bars
add_shape(slide, Inches(0.8), Inches(6.5), Inches(3.8), Inches(0.5), VERDE)
add_textbox(slide, Inches(0.8), Inches(6.5), Inches(3.8), Inches(0.5), '7 treinamento/acompanhamento', font_size=12, bold=True, color=BRANCO, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(4.9), Inches(6.5), Inches(3.8), Inches(0.5), AZUL_MEDIO)
add_textbox(slide, Inches(4.9), Inches(6.5), Inches(3.8), Inches(0.5), '5 em processo', font_size=12, bold=True, color=BRANCO, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(9), Inches(6.5), Inches(3.5), Inches(0.5), CINZA_TEXTO)
add_textbox(slide, Inches(9), Inches(6.5), Inches(3.5), Inches(0.5), '3 aguardando cliente', font_size=12, bold=True, color=BRANCO, alignment=PP_ALIGN.CENTER)

make_slide_number(slide, 6)

# ===================== SLIDE 7: SERVICOS =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CINZA_FUNDO)
make_header_bar(slide, 'SERVICOS DE COMUNICACAO')

make_king_number(slide, 'R$ 1.950', 'total em servicos', Inches(4.5), Inches(1.6), Inches(4), Inches(1.5))

servicos_items = [
    ('1893', 'Troca de titularidade', 'HOFFMAN & PASSOS → AMARAL FARMA ALTO LAGE', 'R$ 550'),
    ('1608', 'Troca de titularidade', 'FAUST → FARMACIA ECONOMICA', 'R$ 550'),
    ('—', 'Inclusao de loja', 'DROGALIMA — incluir loja na comunicacao', 'R$ 850'),
]

for i, (cod, tipo, desc, valor) in enumerate(servicos_items):
    y = 3.0 + i * 1.3
    add_shape(slide, Inches(0.8), Inches(y), Inches(11.7), Inches(1.1), BRANCO)
    add_shape(slide, Inches(0.8), Inches(y), Inches(0.15), Inches(1.1), AZUL_MEDIO)
    add_textbox(slide, Inches(1.2), Inches(y + 0.1), Inches(1.5), Inches(0.4), cod, font_size=14, bold=True, color=AZUL_MEDIO)
    add_textbox(slide, Inches(2.5), Inches(y + 0.1), Inches(3), Inches(0.4), tipo, font_size=16, bold=True, color=AZUL_ESCURO)
    add_textbox(slide, Inches(1.2), Inches(y + 0.5), Inches(8), Inches(0.4), desc, font_size=14, color=CINZA_TEXTO)
    add_textbox(slide, Inches(9.5), Inches(y + 0.2), Inches(2.5), Inches(0.6), valor, font_size=22, bold=True, color=VERDE, alignment=PP_ALIGN.RIGHT)

make_slide_number(slide, 7)

# ===================== SLIDE 8: CLIENTES PERDIDOS =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRANCO)
make_header_bar(slide, 'CLIENTES PERDIDOS — JAN A ABR/2026')

make_king_number(slide, '17', 'clientes perdidos', Inches(1.5), Inches(1.6), Inches(3.5), Inches(1.5))
make_king_number(slide, '-R$ 5.700/mes', 'receita recorrente', Inches(5), Inches(1.6), Inches(4.5), Inches(1.5))
make_king_number(slide, '12', 'so em Abril', Inches(9.5), Inches(1.6), Inches(3), Inches(1.5))

add_shape(slide, Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.04), VERMELHO)

add_textbox(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.4), '3 PRINCIPAIS MOTIVOS', font_size=18, bold=True, color=AZUL_ESCURO)

motivos = [
    ('Troca de sistema / novas funcionalidades', '8 clientes', 'Concorrentes com app web/IA, sistema especializado'),
    ('Falta de atendimento / suporte adequado', '2 clientes', 'Relatos criticos de suporte tecnico inadequado'),
    ('Empresa encerrou atividades', '3 clientes', 'Deram baixa no CNPJ ou encerraram operacoes'),
]

for i, (title, count, desc) in enumerate(motivos):
    y = 4.0 + i * 1.0
    bullet = add_shape(slide, Inches(1.0), Inches(y + 0.05), Inches(0.35), Inches(0.35), VERMELHO if i == 0 else (AMARELO if i == 1 else AZUL_MEDIO))
    add_textbox(slide, Inches(1.6), Inches(y), Inches(4), Inches(0.4), title, font_size=16, bold=True, color=AZUL_ESCURO)
    add_textbox(slide, Inches(6), Inches(y), Inches(2), Inches(0.4), count, font_size=16, bold=True, color=VERMELHO)
    add_textbox(slide, Inches(1.6), Inches(y + 0.4), Inches(9), Inches(0.3), desc, font_size=13, color=CINZA_TEXTO)

add_shape(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.04), AMARELO)
make_slide_number(slide, 8)

# ===================== SLIDE 9: ENCAMINHAMENTOS =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CINZA_FUNDO)
make_header_bar(slide, 'ENCAMINHAMENTOS')

headers = ['ACAO', 'RESPONSAVEL', 'PRAZO']
rows = [
    ['Ciclo completo de atendimento (inicio, meio e fim)', 'Equipe Tecnica', 'Imediato'],
    ['Reforcar proatividade no contato com clientes', 'Todos os tecnicos', 'Continuo'],
    ['Atingir meta de 100 avaliacoes Google', 'Comercial + Suporte', '31/05'],
    ['Acompanhar implantacoes Farma Plus em GO/RO', 'Suporte', 'Continuo'],
    ['Criar plano de retencao para clientes insatisfeitos', 'Supervisao', 'Junho/2026'],
    ['Realizar follow-up pos-implantacao para evitar churn', 'Suporte', 'Imediato'],
]

for j, h in enumerate(headers):
    x = 0.8 + j * 4.0
    add_shape(slide, Inches(x), Inches(1.8), Inches(3.8), Inches(0.6), AZUL_ESCURO)
    add_textbox(slide, Inches(x), Inches(1.8), Inches(3.8), Inches(0.6), h, font_size=14, bold=True, color=BRANCO, alignment=PP_ALIGN.CENTER)

for i, row in enumerate(rows):
    y = 2.5 + i * 0.75
    bg_color = BRANCO if i % 2 == 0 else CINZA_FUNDO
    for j, cell in enumerate(row):
        x = 0.8 + j * 4.0
        add_shape(slide, Inches(x), Inches(y), Inches(3.8), Inches(0.65), bg_color)
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.05), Inches(3.4), Inches(0.55), cell, font_size=13, color=AZUL_ESCURO if j == 0 else CINZA_TEXTO, bold=(j == 0))

make_slide_number(slide, 9)

# ===================== SLIDE 10: ENCERRAMENTO =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, AZUL_ESCURO)
add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2), 'OBRIGADO!', font_size=48, bold=True, color=BRANCO, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8), 'Ate a proxima reuniao', font_size=24, color=AZUL_CLARO, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.5), Inches(5.0), Inches(2.333), Inches(0.06), AZUL_CLARO)
add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6), 'ProSystem Sistemas', font_size=18, color=BRANCO, alignment=PP_ALIGN.CENTER)

# Save
output_path = "squads/administrativo/reunioes/alinhamento/reuniao-alinhamento/output/v1/presentation/Reuniao_Alinhamento_Maio_2026_v3.pptx"
prs.save(output_path)
print(f"OK - Apresentacao salva em: {output_path}")
